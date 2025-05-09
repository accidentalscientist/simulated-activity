import random, subprocess, os, uuid, csv
import pandas as pd
import numpy as np
from datetime import date, datetime, timedelta
from pathlib import Path
from random import choice
from templates import TEMPLATES

# --- parameters -------------------------------------------------------------
NUM_SPRINTS          = 3
COMMITS_PER_SPRINT   = 30
SPRINT_LENGTH_DAYS   = 8
START, END           = date(2022,1,1), date(2025,3,31)
# ---------------------------------------------------------------------------

def create_csv(path):
    path.parent.mkdir(parents=True, exist_ok=True)
    df = pd.DataFrame({
        'feature1': np.random.randn(50),
        'feature2': np.random.randint(0, 100, 50),
        'target'  : np.random.randint(0, 2, 50)
    })
    df.to_csv(path, index=False)

def create_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Sprint update"
        prs.save(path)
    except ImportError:
        # skip PPT if python-pptx not available
        return False
    return True

def create_random_file():
    typ = choice(["py","sql","json","md","csv","pptx"])
    dir_map = {"py":"src","sql":"queries","json":"configs","md":"reports","csv":"data","pptx":"slides"}
    folder = Path(dir_map[typ])
    folder.mkdir(parents=True, exist_ok=True)
    fname = folder / f"{uuid.uuid4().hex[:8]}.{typ}"

    if typ in TEMPLATES:
        fname.write_text(TEMPLATES[typ])
    elif typ == "csv":
        create_csv(fname)
    elif typ == "pptx":
        if not create_pptx(fname):
            return create_random_file()        # fallback to another type
    return fname

def random_sprints():
    span = (END - START).days - (SPRINT_LENGTH_DAYS-1)
    for _ in range(NUM_SPRINTS):
        first = START + timedelta(days=random.randint(0, span))
        for d in range(SPRINT_LENGTH_DAYS):
            day = first + timedelta(days=d)
            daily = COMMITS_PER_SPRINT//SPRINT_LENGTH_DAYS + random.randint(-1,1)
            for _ in range(max(1, daily)):
                sec = random.randint(9*3600, 22*3600)       # 09:00-22:00
                yield datetime.combine(day, datetime.min.time()) + timedelta(seconds=sec)

def make_commit(ts: datetime):
    path = create_random_file()
    Path("dummy.txt").write_text(ts.isoformat())            # touch file
    env = {**os.environ,
           "GIT_AUTHOR_DATE": ts.isoformat(),
           "GIT_COMMITTER_DATE": ts.isoformat()}
    subprocess.run(["git","add","dummy.txt"], check=True)
    subprocess.run(["git","commit","-m", f"sim {ts:%F %T}"], check=True, env=env)

if __name__ == "__main__":
    for t in sorted(random_sprints()):
        make_commit(t)
